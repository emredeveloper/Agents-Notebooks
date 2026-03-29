from fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path


mcp = FastMCP("powerpoint")


def ensure_parent_dir(file_path: str) -> None:
    Path(file_path).parent.mkdir(parents=True, exist_ok=True)


def ensure_presentation(file_path: str) -> Presentation:
    ensure_parent_dir(file_path)
    if Path(file_path).exists():
        return Presentation(file_path)

    prs = Presentation()
    prs.save(file_path)
    return Presentation(file_path)


def save_presentation(prs: Presentation, file_path: str) -> str:
    prs.save(file_path)
    return f"Kaydedildi: {file_path}"


def parse_rgb(color: str | None) -> RGBColor | None:
    if not color:
        return None

    value = color.strip().replace("#", "")
    if len(value) != 6:
        raise ValueError("Renk 6 haneli hex olmalı. Örnek: 'FF0000' veya '#FF0000'")

    return RGBColor(
        int(value[0:2], 16),
        int(value[2:4], 16),
        int(value[4:6], 16),
    )


def parse_alignment(alignment: str | None):
    if not alignment:
        return PP_ALIGN.LEFT

    value = alignment.strip().lower()
    mapping = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    return mapping.get(value, PP_ALIGN.LEFT)


def set_run_font(run, font_name, font_size, bold, italic, color):
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    rgb = parse_rgb(color)
    if rgb:
        font.color.rgb = rgb


def style_paragraph(
    paragraph,
    text: str,
    font_name: str = "Calibri",
    font_size: int = 20,
    bold: bool = False,
    italic: bool = False,
    color: str | None = None,
    alignment: str = "left",
):
    paragraph.clear()
    paragraph.alignment = parse_alignment(alignment)
    run = paragraph.add_run()
    run.text = text
    set_run_font(run, font_name, font_size, bold, italic, color)


def style_text_frame_all_paragraphs(
    text_frame,
    font_name: str = "Calibri",
    font_size: int = 20,
    bold: bool = False,
    italic: bool = False,
    color: str | None = None,
    alignment: str = "left",
):
    for p in text_frame.paragraphs:
        p.alignment = parse_alignment(alignment)
        if not p.runs:
            run = p.add_run()
            run.text = p.text or ""
        for run in p.runs:
            set_run_font(run, font_name, font_size, bold, italic, color)


@mcp.tool()
def create_presentation(file_path: str) -> str:
    prs = ensure_presentation(file_path)
    prs.save(file_path)
    return f"Sunum hazır: {file_path}"


@mcp.tool()
def get_slide_count(file_path: str) -> str:
    prs = ensure_presentation(file_path)
    return f"Toplam slayt sayısı: {len(prs.slides)}"


@mcp.tool()
def add_title_slide(
    file_path: str,
    title: str,
    subtitle: str = "",
    title_font_name: str = "Calibri",
    title_font_size: int = 30,
    title_bold: bool = True,
    title_italic: bool = False,
    title_color: str = "000000",
    subtitle_font_name: str = "Calibri",
    subtitle_font_size: int = 18,
    subtitle_bold: bool = False,
    subtitle_italic: bool = False,
    subtitle_color: str = "444444",
    title_alignment: str = "center",
    subtitle_alignment: str = "center",
) -> str:
    prs = ensure_presentation(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    style_text_frame_all_paragraphs(
        title_shape.text_frame,
        font_name=title_font_name,
        font_size=title_font_size,
        bold=title_bold,
        italic=title_italic,
        color=title_color,
        alignment=title_alignment,
    )

    style_text_frame_all_paragraphs(
        subtitle_shape.text_frame,
        font_name=subtitle_font_name,
        font_size=subtitle_font_size,
        bold=subtitle_bold,
        italic=subtitle_italic,
        color=subtitle_color,
        alignment=subtitle_alignment,
    )

    return save_presentation(prs, file_path)


@mcp.tool()
def add_text_slide(
    file_path: str,
    title: str,
    content: str,
    title_font_name: str = "Calibri",
    title_font_size: int = 28,
    title_bold: bool = True,
    title_italic: bool = False,
    title_color: str = "000000",
    title_alignment: str = "left",
    content_font_name: str = "Calibri",
    content_font_size: int = 20,
    content_bold: bool = False,
    content_italic: bool = False,
    content_color: str = "222222",
    content_alignment: str = "left",
) -> str:
    prs = ensure_presentation(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = title
    body_shape.text = content

    style_text_frame_all_paragraphs(
        title_shape.text_frame,
        font_name=title_font_name,
        font_size=title_font_size,
        bold=title_bold,
        italic=title_italic,
        color=title_color,
        alignment=title_alignment,
    )

    style_text_frame_all_paragraphs(
        body_shape.text_frame,
        font_name=content_font_name,
        font_size=content_font_size,
        bold=content_bold,
        italic=content_italic,
        color=content_color,
        alignment=content_alignment,
    )

    return save_presentation(prs, file_path)


@mcp.tool()
def add_bullets_slide(
    file_path: str,
    title: str,
    bullets: list[str],
    title_font_name: str = "Calibri",
    title_font_size: int = 28,
    title_bold: bool = True,
    title_italic: bool = False,
    title_color: str = "000000",
    title_alignment: str = "left",
    bullet_font_name: str = "Calibri",
    bullet_font_size: int = 20,
    bullet_bold: bool = False,
    bullet_italic: bool = False,
    bullet_color: str = "222222",
    bullet_alignment: str = "left",
) -> str:
    prs = ensure_presentation(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = title
    style_text_frame_all_paragraphs(
        title_shape.text_frame,
        font_name=title_font_name,
        font_size=title_font_size,
        bold=title_bold,
        italic=title_italic,
        color=title_color,
        alignment=title_alignment,
    )

    tf = body_shape.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        style_paragraph(
            p,
            text=bullet,
            font_name=bullet_font_name,
            font_size=bullet_font_size,
            bold=bullet_bold,
            italic=bullet_italic,
            color=bullet_color,
            alignment=bullet_alignment,
        )

    return save_presentation(prs, file_path)


@mcp.tool()
def add_textbox_to_blank_slide(
    file_path: str,
    text: str,
    left: float = 1.0,
    top: float = 1.0,
    width: float = 8.0,
    height: float = 2.0,
    font_name: str = "Calibri",
    font_size: int = 24,
    bold: bool = False,
    italic: bool = False,
    color: str = "000000",
    alignment: str = "left",
) -> str:
    prs = ensure_presentation(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    textbox = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )

    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]

    style_paragraph(
        p,
        text=text,
        font_name=font_name,
        font_size=font_size,
        bold=bold,
        italic=italic,
        color=color,
        alignment=alignment,
    )

    return save_presentation(prs, file_path)


@mcp.tool()
def add_title_and_textboxes_slide(
    file_path: str,
    title: str,
    box1_text: str,
    box2_text: str = "",
    title_font_size: int = 28,
    body_font_size: int = 20,
    font_name: str = "Calibri",
    title_color: str = "000000",
    body_color: str = "222222",
) -> str:
    prs = ensure_presentation(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    style_paragraph(
        p,
        text=title,
        font_name=font_name,
        font_size=title_font_size,
        bold=True,
        italic=False,
        color=title_color,
        alignment="left",
    )

    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.5))
    p1 = left_box.text_frame.paragraphs[0]
    style_paragraph(
        p1,
        text=box1_text,
        font_name=font_name,
        font_size=body_font_size,
        bold=False,
        italic=False,
        color=body_color,
        alignment="left",
    )

    if box2_text.strip():
        right_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.6), Inches(5.5), Inches(4.5))
        p2 = right_box.text_frame.paragraphs[0]
        style_paragraph(
            p2,
            text=box2_text,
            font_name=font_name,
            font_size=body_font_size,
            bold=False,
            italic=False,
            color=body_color,
            alignment="left",
        )

    return save_presentation(prs, file_path)


@mcp.tool()
def add_table_slide(
    file_path: str,
    title: str,
    data: list[list[str]],
    left: float = 0.7,
    top: float = 1.6,
    width: float = 11.5,
    height: float = 4.5,
    title_font_name: str = "Calibri",
    title_font_size: int = 28,
    title_bold: bool = True,
    title_italic: bool = False,
    title_color: str = "000000",
    title_alignment: str = "left",
    header_font_name: str = "Calibri",
    header_font_size: int = 18,
    header_bold: bool = True,
    header_italic: bool = False,
    header_color: str = "000000",
    header_alignment: str = "center",
    body_font_name: str = "Calibri",
    body_font_size: int = 16,
    body_bold: bool = False,
    body_italic: bool = False,
    body_color: str = "222222",
    body_alignment: str = "center",
) -> str:
    """
    data parametresi 2 boyutlu liste olmalı.
    İlk satır header kabul edilir.
    Örnek:
    [
      ["Ürün", "Fiyat", "Adet"],
      ["Elma", "10", "5"],
      ["Armut", "12", "3"]
    ]
    """
    if not data or not isinstance(data, list) or not all(isinstance(row, list) for row in data):
        return "Hata: data, 2 boyutlu liste olmalı."

    row_count = len(data)
    col_count = max(len(row) for row in data)
    if row_count == 0 or col_count == 0:
        return "Hata: tablo verisi boş."

    prs = ensure_presentation(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    style_paragraph(
        p,
        text=title,
        font_name=title_font_name,
        font_size=title_font_size,
        bold=title_bold,
        italic=title_italic,
        color=title_color,
        alignment=title_alignment,
    )

    table_shape = slide.shapes.add_table(
        row_count,
        col_count,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    table = table_shape.table

    for r in range(row_count):
        for c in range(col_count):
            value = data[r][c] if c < len(data[r]) else ""
            cell = table.cell(r, c)
            cell.text = str(value)

            tf = cell.text_frame
            if not tf.paragraphs:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]

            if r == 0:
                style_paragraph(
                    p,
                    text=str(value),
                    font_name=header_font_name,
                    font_size=header_font_size,
                    bold=header_bold,
                    italic=header_italic,
                    color=header_color,
                    alignment=header_alignment,
                )
            else:
                style_paragraph(
                    p,
                    text=str(value),
                    font_name=body_font_name,
                    font_size=body_font_size,
                    bold=body_bold,
                    italic=body_italic,
                    color=body_color,
                    alignment=body_alignment,
                )

    return save_presentation(prs, file_path)


@mcp.tool()
def style_last_slide_title(
    file_path: str,
    font_name: str = "Calibri",
    font_size: int = 28,
    bold: bool = True,
    italic: bool = False,
    color: str = "000000",
    alignment: str = "left",
) -> str:
    prs = ensure_presentation(file_path)

    if len(prs.slides) == 0:
        return "Hata: Sunumda hiç slayt yok."

    slide = prs.slides[-1]
    if not slide.shapes.title:
        return "Hata: Son slaytta başlık alanı bulunamadı."

    style_text_frame_all_paragraphs(
        slide.shapes.title.text_frame,
        font_name=font_name,
        font_size=font_size,
        bold=bold,
        italic=italic,
        color=color,
        alignment=alignment,
    )

    return save_presentation(prs, file_path)


@mcp.tool()
def style_last_slide_body(
    file_path: str,
    font_name: str = "Calibri",
    font_size: int = 20,
    bold: bool = False,
    italic: bool = False,
    color: str = "222222",
    alignment: str = "left",
) -> str:
    prs = ensure_presentation(file_path)

    if len(prs.slides) == 0:
        return "Hata: Sunumda hiç slayt yok."

    slide = prs.slides[-1]

    body_shape = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape != slide.shapes.title:
            body_shape = shape
            break

    if body_shape is None:
        return "Hata: Son slaytta düzenlenecek metin alanı bulunamadı."

    style_text_frame_all_paragraphs(
        body_shape.text_frame,
        font_name=font_name,
        font_size=font_size,
        bold=bold,
        italic=italic,
        color=color,
        alignment=alignment,
    )

    return save_presentation(prs, file_path)


if __name__ == "__main__":
    mcp.run()